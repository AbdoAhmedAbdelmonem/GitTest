import os
import json
import logging
from typing import List, Dict, Any, Optional
import openai
import google.generativeai as genai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
import PyPDF2
from pptx import Presentation
import nbformat
import uuid
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        self.embeddings = OpenAIEmbeddings()
        self.vector_store = None
        self.documents = {}
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def extract_text_from_notebook(self, file_path: str) -> str:
        """Extract text from Jupyter notebook"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                nb = nbformat.read(f, as_version=4)
            
            text = ""
            for cell in nb.cells:
                if cell.cell_type == 'markdown':
                    text += cell.source + "\n\n"
                elif cell.cell_type == 'code':
                    text += f"```python\n{cell.source}\n```\n\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting notebook text: {e}")
            return ""
    
    def extract_text_from_python(self, file_path: str) -> str:
        """Extract text from Python file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading Python file: {e}")
            return ""
    
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """Process a document and add it to the vector store"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Extract text based on file type
        if file_extension == '.pdf':
            text = self.extract_text_from_pdf(file_path)
        elif file_extension == '.pptx':
            text = self.extract_text_from_pptx(file_path)
        elif file_extension == '.ipynb':
            text = self.extract_text_from_notebook(file_path)
        elif file_extension == '.py':
            text = self.extract_text_from_python(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        if not text.strip():
            raise ValueError("No text could be extracted from the document")
        
        # Split text into chunks
        chunks = self.text_splitter.split_text(text)
        
        # Generate document ID
        doc_id = str(uuid.uuid4())
        
        # Store document metadata
        self.documents[doc_id] = {
            "filename": os.path.basename(file_path),
            "file_type": file_extension,
            "processed_at": datetime.now().isoformat(),
            "chunks_count": len(chunks)
        }
        
        # Create or update vector store
        if self.vector_store is None:
            self.vector_store = FAISS.from_texts(chunks, self.embeddings)
        else:
            self.vector_store.add_texts(chunks)
        
        return {
            "document_id": doc_id,
            "chunks": chunks,
            "metadata": self.documents[doc_id]
        }
    
    def get_document_list(self) -> List[Dict[str, Any]]:
        """Get list of all processed documents"""
        return list(self.documents.values())
    
    def delete_document(self, doc_id: str):
        """Delete a document from the system"""
        if doc_id in self.documents:
            del self.documents[doc_id]
            # Note: FAISS doesn't support individual document deletion easily
            # In production, you'd need to rebuild the index
        else:
            raise ValueError(f"Document {doc_id} not found")

class ChatBot:
    def __init__(self):
        self.processor = DocumentProcessor()
        
        # Initialize AI models
        self.openai_client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
        self.google_model = genai.GenerativeModel('gemini-pro')
    
    def get_response(self, question: str, language: str = "arabic", model: str = "openai") -> Dict[str, Any]:
        """Get AI response to a question"""
        try:
            # Get relevant context from documents
            context = ""
            sources = []
            
            if self.processor.vector_store:
                docs = self.processor.vector_store.similarity_search(question, k=3)
                context = "\n".join([doc.page_content for doc in docs])
                sources = [f"Document chunk {i+1}" for i in range(len(docs))]
            
            # Prepare prompt
            if language == "arabic":
                system_prompt = """أنت مساعد ذكي يجيب على الأسئلة باللغة العربية. 
                استخدم المعلومات المتوفرة في السياق للإجابة على السؤال بدقة ووضوح."""
            else:
                system_prompt = """You are an intelligent assistant that answers questions accurately and clearly. 
                Use the provided context to answer the question."""
            
            prompt = f"{system_prompt}\n\nContext: {context}\n\nQuestion: {question}\n\nAnswer:"
            
            # Get response from selected model
            if model == "openai":
                response = self.openai_client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"Context: {context}\n\nQuestion: {question}"}
                    ],
                    max_tokens=500,
                    temperature=0.7
                )
                answer = response.choices[0].message.content
            
            elif model == "google":
                response = self.google_model.generate_content(prompt)
                answer = response.text
            
            else:
                answer = "Model not supported"
            
            return {
                "answer": answer,
                "sources": sources,
                "confidence": 0.8,  # You can implement confidence scoring
            }
        
        except Exception as e:
            logger.error(f"Error generating response: {e}")
            return {
                "answer": "عذراً، حدث خطأ في معالجة سؤالك. يرجى المحاولة مرة أخرى.",
                "sources": [],
                "confidence": 0.0,
                "error": str(e)
            }
